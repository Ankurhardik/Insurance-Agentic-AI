from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, BackgroundTasks
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from sqlalchemy import select
import os
import uuid
import io
import json
import csv

from Backend.src.db.session import get_db
from Backend.src.models.document import KnowledgeDocument
from Backend.src.schemas.document import DocumentOut, DocumentDetailOut
from Backend.src.api.endpoints.auth import get_current_user, RoleChecker
from Backend.src.models.user import User

router = APIRouter(prefix="/api/documents", tags=["documents"])

# ─── Allowed file types ──────────────────────────────────────────────────────
ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".xlsx", ".xls", ".csv", ".json",
    ".txt", ".md", ".yaml", ".yml",
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
}

ALLOWED_CONTENT_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/csv",
    "application/json",
    "text/plain",
    "text/markdown",
    "application/x-yaml",
    "text/yaml",
    "image/png",
    "image/jpeg",
    "image/gif",
    "image/webp",
    "image/svg+xml",
    # Browsers sometimes use generic for uploads — we rely on extension too
    "application/octet-stream",
}


def extract_text(content: bytes, ext: str, filename: str) -> str:
    """Extract text content from various document formats."""
    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"--- Page {i} ---\n{text}")
                return "\n\n".join(pages) if pages else "[No text could be extracted from this PDF]"

        elif ext in (".docx",):
            from docx import Document
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n".join(paragraphs) if paragraphs else "[No text content found in this document]"

        elif ext in (".pptx",):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            return "\n\n".join(slides_text) if slides_text else "[No text content found in this presentation]"

        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(sheets_text) if sheets_text else "[No data found in this spreadsheet]"

        elif ext == ".csv":
            decoded = content.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(decoded))
            rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
            return "\n".join(rows) if rows else "[Empty CSV file]"

        elif ext == ".json":
            decoded = content.decode("utf-8", errors="replace")
            parsed = json.loads(decoded)
            return json.dumps(parsed, indent=2, ensure_ascii=False)

        elif ext in (".txt", ".md", ".yaml", ".yml"):
            return content.decode("utf-8", errors="replace")

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"):
            return "[Image file - Preview visual directly]"

        else:
            return f"[Preview not supported for '{ext}' file format]"

    except Exception as e:
        return f"[Error extracting content: {str(e)}]"


def process_document_embeddings(doc_id: str):
    """Asynchronously split document text and generate vector embeddings using BAAI/bge-large-en-v1.5."""
    from Backend.src.db.session import SessionLocal, VectorSessionLocal
    from Backend.src.models.document import KnowledgeDocument
    from Backend.src.models.chunk import DocumentChunk
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from sentence_transformers import SentenceTransformer
    import uuid

    db = SessionLocal()
    vector_db = VectorSessionLocal()
    try:
        # Get the document
        doc = db.query(KnowledgeDocument).filter(KnowledgeDocument.id == doc_id).first()
        if not doc:
            print(f"[Embeddings] Document {doc_id} not found.")
            return

        # Update status to Processing
        doc.status = "Processing"
        db.commit()

        text = doc.extracted_text
        # If it's an image or has no text, skip embedding generation
        if not text or text.strip() == "" or text.startswith("[Image file") or text.startswith("[Preview not supported") or text.startswith("[Error extracting"):
            print(f"[Embeddings] Skipping embedding generation for {doc.filename} (invalid/empty text content).")
            doc.status = "Ready"
            db.commit()
            return

        # Split text into chunks
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=50,
            length_function=len
        )
        chunks = splitter.split_text(text)
        print(f"[Embeddings] Split {doc.filename} into {len(chunks)} chunks.")

        if chunks:
            # Load sentence-transformers model BAAI/bge-large-en-v1.5
            print(f"[Embeddings] Loading BAAI/bge-large-en-v1.5 model...")
            model = SentenceTransformer('BAAI/bge-large-en-v1.5')
            
            print(f"[Embeddings] Encoding {len(chunks)} chunks...")
            embeddings = model.encode(chunks, show_progress_bar=True)

            # Clear existing chunks for this document in pgvector database
            vector_db.query(DocumentChunk).filter(DocumentChunk.document_id == doc_id).delete()

            # Insert new chunks
            for idx, (chunk_text, emb) in enumerate(zip(chunks, embeddings)):
                chunk_record = DocumentChunk(
                    id=str(uuid.uuid4()),
                    document_id=doc_id,
                    chunk_index=idx,
                    content=chunk_text,
                    embedding=list(emb)  # Convert numpy array to list
                )
                vector_db.add(chunk_record)

            vector_db.commit()
            print(f"[Embeddings] Saved {len(chunks)} chunks to vector database for {doc.filename}.")

        doc.status = "Ready"
        db.commit()
    except Exception as e:
        print(f"[Embeddings] Error generating embeddings for document {doc_id}: {e}")
        try:
            doc.status = "Failed"
            db.commit()
        except Exception:
            pass
        vector_db.rollback()
    finally:
        db.close()
        vector_db.close()


# ─── Routes ──────────────────────────────────────────────────────────────────

@router.get("", response_model=list[DocumentOut])
def list_documents(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """List all documents (authenticated users)."""
    stmt = select(KnowledgeDocument).order_by(KnowledgeDocument.created_at.desc())
    result = db.execute(stmt)
    return result.scalars().all()


@router.get("/serve/{doc_id}")
def serve_document(
    doc_id: str,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Stream the raw file for in-browser viewing (PDF embed, etc.)."""
    stmt = select(KnowledgeDocument).where(KnowledgeDocument.id == doc_id)
    doc = db.execute(stmt).scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found.")
    if not doc.file_path or not os.path.exists(doc.file_path):
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found on disk.")
    return FileResponse(
        path=doc.file_path,
        media_type=doc.content_type,
        filename=doc.filename,
        headers={"Content-Disposition": f'inline; filename="{doc.filename}"'}
    )


@router.get("/{doc_id}", response_model=DocumentDetailOut)
def get_document(
    doc_id: str,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get single document details including extracted text (authenticated users)."""
    stmt = select(KnowledgeDocument).where(KnowledgeDocument.id == doc_id)
    doc = db.execute(stmt).scalar_one_or_none()
    if not doc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found."
        )
    return doc


@router.post("/upload", response_model=list[DocumentOut], status_code=status.HTTP_201_CREATED)
async def upload_documents(
    background_tasks: BackgroundTasks,
    files: list[UploadFile] = File(...),
    current_user: User = Depends(RoleChecker(["admin"])),
    db: Session = Depends(get_db)
):
    """Multi-upload documents (Admin only). Supports PDF, DOCX, PPTX, XLSX, CSV, JSON, TXT, MD."""
    upload_dir = os.path.join("Backend", "uploads")
    os.makedirs(upload_dir, exist_ok=True)

    docs_created = []

    for file in files:
        filename = file.filename or "unnamed_file"
        ext = os.path.splitext(filename)[1].lower()

        # Validate extension
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"File type '{ext}' is not supported. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
            )

        doc_id = str(uuid.uuid4())
        saved_filename = f"{doc_id}_{filename}"
        file_path = os.path.join(upload_dir, saved_filename)

        content = await file.read()

        with open(file_path, "wb") as f:
            f.write(content)

        extracted_text = extract_text(content, ext, filename)

        db_doc = KnowledgeDocument(
            id=doc_id,
            filename=filename,
            content_type=file.content_type or "application/octet-stream",
            file_size=len(content),
            status="Pending",
            file_path=file_path,
            extracted_text=extracted_text,
            created_by=current_user.id
        )
        db.add(db_doc)
        docs_created.append(db_doc)

    db.commit()
    for doc in docs_created:
        db.refresh(doc)
        background_tasks.add_task(process_document_embeddings, doc.id)

    return docs_created


@router.delete("/{doc_id}")
def delete_document(
    doc_id: str,
    current_user: User = Depends(RoleChecker(["admin"])),
    db: Session = Depends(get_db)
):
    """Delete a document (Admin only)."""
    stmt = select(KnowledgeDocument).where(KnowledgeDocument.id == doc_id)
    doc = db.execute(stmt).scalar_one_or_none()
    if not doc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Document not found."
        )

    # Delete corresponding chunks from vector database
    from Backend.src.db.session import VectorSessionLocal
    from Backend.src.models.chunk import DocumentChunk
    vector_db = VectorSessionLocal()
    try:
        vector_db.query(DocumentChunk).filter(DocumentChunk.document_id == doc_id).delete()
        vector_db.commit()
    except Exception as e:
        print(f"[Embeddings] Failed to delete chunks for document {doc_id} on deletion: {e}")
        vector_db.rollback()
    finally:
        vector_db.close()

    if doc.file_path and os.path.exists(doc.file_path):
        try:
            os.remove(doc.file_path)
        except Exception:
            pass

    db.delete(doc)
    db.commit()
    return {"detail": "Document successfully deleted."}
